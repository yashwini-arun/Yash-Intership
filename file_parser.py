import io
from typing import Tuple


def extract_text(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    """
    Extract plain text from any file format.
    Returns (extracted_text, detected_format)
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    # ── PDF ──────────────────────────────────────────────────────────────────
    if ext == "pdf":
        return _extract_pdf(file_bytes), "pdf"

    # ── Word Document ─────────────────────────────────────────────────────────
    if ext in ("docx", "doc"):
        return _extract_docx(file_bytes), "docx"

    # ── PowerPoint ────────────────────────────────────────────────────────────
    if ext in ("pptx", "ppt"):
        return _extract_pptx(file_bytes), "pptx"

    # ── Excel ─────────────────────────────────────────────────────────────────
    if ext in ("xlsx", "xls"):
        return _extract_excel(file_bytes), "xlsx"

    # ── CSV ───────────────────────────────────────────────────────────────────
    if ext == "csv":
        return file_bytes.decode("utf-8", errors="ignore"), "csv"

    # ── JSON ──────────────────────────────────────────────────────────────────
    if ext == "json":
        import json
        try:
            data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
            return json.dumps(data, indent=2), "json"
        except:
            return file_bytes.decode("utf-8", errors="ignore"), "json"

    # ── Markdown / TXT / any text file ───────────────────────────────────────
    return file_bytes.decode("utf-8", errors="ignore"), ext or "txt"


# ── Extractors ────────────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> str:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i+1}]\n{text.strip()}")
        return "\n\n".join(pages)
    except ImportError:
        raise ImportError("PyPDF2 not installed. Run: pip install PyPDF2")
    except Exception as e:
        raise ValueError(f"Could not read PDF: {str(e)}")


def _extract_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")
    except Exception as e:
        raise ValueError(f"Could not read DOCX: {str(e)}")


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"Could not read PPTX: {str(e)}")


def _extract_excel(file_bytes: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")
    except Exception as e:
        raise ValueError(f"Could not read Excel: {str(e)}")